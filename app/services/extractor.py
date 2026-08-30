import os
import re
import math
import base64
from io import BytesIO
from PIL import Image, ImageDraw, ImageFont
import fitz
from app.config import settings

TEMP_GRID_DIR = os.path.join("uploads", "temp_grids")
os.makedirs(TEMP_GRID_DIR, exist_ok=True)


def format_extracted_text(raw_text: str) -> str:
    if not raw_text or not raw_text.strip():
        return raw_text

    text = raw_text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    lines = [l.strip() for l in text.split("\n") if l.strip()]

    formatted = []
    for l in lines:
        if re.match(r"^--- Page \d+", l, re.I) or re.match(r"^--- Slide \d+", l, re.I):
            formatted.append(f"\n\n---\n\n### {l.strip('- ')}\n")
        elif re.match(r"^(subject\s*[:.]|question\s*bank|unit\s+\d+|chapter\s+\d+|module\s+\d+|topic\s*[:.])", l, re.I):
            formatted.append(f"\n## {l}\n")
        elif re.match(r"^(short\s+questions?|broad\s+questions?|long\s+questions?|section\s+[a-z]|fill\s+in|true\s+or\s+false|match\s+the|multiple\s+choice|mcq|objective)", l, re.I):
            formatted.append(f"\n### {l}\n")
        elif re.match(r"^(Q\.?\s*)?(\d{1,3})\s*[.)]\s+", l):
            formatted.append(f"\n**{l}**\n")
        elif l.startswith(("-", "*", "•")):
            formatted.append(f"- {l.lstrip('-*• ')}")
        elif re.match(r"^\(?[a-eA-Eivx]{1,4}[.)]\s+", l):
            formatted.append(f"   - {l}")
        else:
            formatted.append(l)

    result = "\n".join(formatted).strip()
    return re.sub(r"\n{3,}", "\n\n", result)


def stitch_pages_to_2x2_grid(page_images: list[Image.Image], page_numbers: list[int], output_path: str) -> str:
    num_pages = len(page_images)
    if num_pages == 0:
        return ""

    cols = 2 if num_pages > 1 else 1
    rows = 2 if num_pages > 2 else 1

    cell_w, cell_h = 750, 1000
    label_h = 40
    gap = 12

    grid_w = (cols * cell_w) + ((cols + 1) * gap)
    grid_h = (rows * (cell_h + label_h)) + ((rows + 1) * gap)

    grid_img = Image.new("RGB", (grid_w, grid_h), color=(15, 23, 42))
    draw = ImageDraw.Draw(grid_img)

    try:
        font = ImageFont.truetype("arial.ttf", 22)
    except Exception:
        font = ImageFont.load_default()

    for idx, (img, p_num) in enumerate(zip(page_images, page_numbers)):
        r = idx // cols
        c = idx % cols

        x = gap + c * (cell_w + gap)
        y = gap + r * (cell_h + label_h + gap)

        img_resized = img.copy()
        img_resized.thumbnail((cell_w, cell_h), Image.Resampling.LANCZOS)

        pad = Image.new("RGB", (cell_w, cell_h), color=(255, 255, 255))
        ox = (cell_w - img_resized.width) // 2
        oy = (cell_h - img_resized.height) // 2
        pad.paste(img_resized, (ox, oy))

        grid_img.paste(pad, (x, y))

        label_text = f"=== PAGE {p_num} ==="
        bbox = draw.textbbox((0, 0), label_text, font=font)
        text_w = bbox[2] - bbox[0]
        text_x = x + (cell_w - text_w) // 2
        text_y = y + cell_h + 8
        draw.text((text_x, text_y), label_text, fill=(224, 231, 255), font=font)

    grid_img.save(output_path, format="JPEG", quality=85, optimize=True)
    return output_path


def scan_grid_image_with_vision(grid_path: str, page_numbers: list[int]) -> str:
    import time
    from google import genai

    if not os.path.exists(grid_path):
        return ""

    pages_str = ", ".join(f"Page {p}" for p in page_numbers)
    prompt = (
        f"Transcribe all text, formulas (using LaTeX $...$), and tables from this 2x2 grid containing {pages_str}.\n"
        "Format output strictly by page headers:\n"
        f"--- Page {page_numbers[0]} ---\n[Extracted Text]\n"
    )

    if not settings.GEMINI_API_KEY:
        return ""

    try:
        client = genai.Client(api_key=settings.GEMINI_API_KEY.strip())
        with open(grid_path, "rb") as f:
            img_bytes = f.read()

        candidate_models = [
            "gemini-3.5-flash-lite",
            "gemini-flash-latest",
            "gemini-3.5-flash",
            "gemini-2.5-flash"
        ]

        for attempt in range(3):
            for model_name in candidate_models:
                try:
                    response = client.models.generate_content(
                        model=model_name,
                        contents=[
                            {
                                "role": "user",
                                "parts": [
                                    {"text": prompt},
                                    {"inline_data": {"mime_type": "image/jpeg", "data": img_bytes}}
                                ]
                            }
                        ]
                    )
                    if response.text and len(response.text.strip()) > 30:
                        return response.text.strip()
                except Exception as e:
                    time.sleep(1.0)
            time.sleep(2.0)
    except Exception as e:
        print(f"Vision 2x2 grid scan error: {e}")

    return ""


def extract_scanned_pdf_via_2x2_grid(file_path: str, total_pages: int, on_page_progress=None) -> str:
    import time
    filename = os.path.basename(file_path)
    extracted_batches = []
    batch_size = 4

    try:
        doc = fitz.open(file_path)
        num_batches = math.ceil(total_pages / batch_size)

        for b_idx in range(num_batches):
            start_p = b_idx * batch_size
            end_p = min(start_p + batch_size, total_pages)

            page_images = []
            page_numbers = []

            for p_idx in range(start_p, end_p):
                page = doc[p_idx]
                pix = page.get_pixmap(dpi=150)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                page_images.append(img)
                page_numbers.append(p_idx + 1)

            grid_filename = f"grid_{os.path.splitext(filename)[0]}_batch_{b_idx + 1}.jpg"
            grid_path = os.path.join(TEMP_GRID_DIR, grid_filename)

            stitch_pages_to_2x2_grid(page_images, page_numbers, grid_path)

            batch_text = scan_grid_image_with_vision(grid_path, page_numbers)
            if batch_text:
                extracted_batches.append(batch_text)

            try:
                if os.path.exists(grid_path):
                    os.remove(grid_path)
            except Exception:
                pass

            if on_page_progress:
                on_page_progress(end_p, total_pages, True)

            time.sleep(1.2)

        doc.close()
    except Exception as e:
        print(f"Error in 2x2 grid scanned PDF extraction: {e}")

    full_text = "\n\n".join(extracted_batches).strip()
    return format_extracted_text(full_text) if full_text else f"Scanned Study Material: '{filename}'"


def extract_pdf_clean(file_path: str, on_page_progress=None) -> str:
    filename = os.path.basename(file_path)
    extracted_pages = []
    has_digital_text = False
    total_pages = 0

    try:
        doc = fitz.open(file_path)
        total_pages = len(doc)
        total_chars = 0

        for page_idx in range(total_pages):
            page_num = page_idx + 1
            page = doc[page_idx]
            page_text = page.get_text("text").strip()

            if page_text:
                total_chars += len(page_text)
                extracted_pages.append(f"--- Page {page_num} of {total_pages} ---\n{page_text}")

            if on_page_progress:
                on_page_progress(page_num, total_pages, False)

        doc.close()
        has_digital_text = (total_chars > (total_pages * 15))
    except Exception as e:
        print(f"Error reading PDF {filename}: {e}")

    if has_digital_text and extracted_pages:
        full_text = "\n\n".join(extracted_pages).strip()
        return format_extracted_text(full_text)

    return extract_scanned_pdf_via_2x2_grid(file_path, total_pages=max(1, total_pages), on_page_progress=on_page_progress)


def extract_docx_clean(file_path: str, on_page_progress=None) -> str:
    filename = os.path.basename(file_path)
    try:
        import docx
        doc = docx.Document(file_path)
        parts = []

        for p in doc.paragraphs:
            txt = p.text.strip()
            if not txt:
                continue
            if p.style.name.startswith("Heading 1"):
                parts.append(f"\n# {txt}\n")
            elif p.style.name.startswith("Heading 2"):
                parts.append(f"\n## {txt}\n")
            elif p.style.name.startswith("Heading 3"):
                parts.append(f"\n### {txt}\n")
            elif p.style.name.startswith("List"):
                parts.append(f"- {txt}")
            else:
                parts.append(txt)

        for t_idx, table in enumerate(doc.tables):
            table_rows = []
            for r_idx, row in enumerate(table.rows):
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                row_str = "| " + " | ".join(cells) + " |"
                table_rows.append(row_str)
                if r_idx == 0:
                    sep = "| " + " | ".join(["---"] * len(cells)) + " |"
                    table_rows.append(sep)
            if table_rows:
                parts.append("\n" + "\n".join(table_rows) + "\n")

        if on_page_progress:
            on_page_progress(1, 1, False)

        text = "\n".join(parts).strip()
        return format_extracted_text(text) if text else f"Word Document Note: '{filename}'"
    except Exception as e:
        print(f"Error extracting DOCX {filename}: {e}")
        return f"Word Document Note: '{filename}'"


def extract_pptx_clean(file_path: str, on_page_progress=None) -> str:
    filename = os.path.basename(file_path)
    try:
        import pptx
        prs = pptx.Presentation(file_path)
        total_slides = len(prs.slides)
        slides_text = []

        for s_idx, slide in enumerate(prs.slides):
            slide_num = s_idx + 1
            shape_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    txt = shape.text.strip()
                    shape_texts.append(txt)

            if shape_texts:
                slides_text.append(f"--- Slide {slide_num} of {total_slides} ---\n" + "\n".join(shape_texts))

            if on_page_progress:
                on_page_progress(slide_num, total_slides, False)

        text = "\n\n".join(slides_text).strip()
        return format_extracted_text(text) if text else f"Presentation Lecture: '{filename}'"
    except Exception as e:
        print(f"Error extracting PPTX {filename}: {e}")
        return f"Presentation Lecture: '{filename}'"


def extract_text_file_clean(file_path: str, on_page_progress=None) -> str:
    filename = os.path.basename(file_path)
    encodings = ["utf-8", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc, errors="ignore") as f:
                text = f.read().strip()
                if text:
                    if on_page_progress:
                        on_page_progress(1, 1, False)
                    return format_extracted_text(text)
        except Exception:
            continue
    return f"Study Note: '{filename}'"


def extract_text_from_file(file_path: str, doc_id: int | None = None, unique_code: str | None = None, classroom_code: str | None = None, on_page_progress=None) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_pdf_clean(file_path, on_page_progress=on_page_progress)

    if ext in (".docx", ".doc"):
        return extract_docx_clean(file_path, on_page_progress=on_page_progress)

    if ext in (".pptx", ".ppt"):
        return extract_pptx_clean(file_path, on_page_progress=on_page_progress)

    if ext in (".txt", ".md", ".json", ".csv", ".py", ".c", ".cpp", ".java"):
        return extract_text_file_clean(file_path, on_page_progress=on_page_progress)

    return extract_text_file_clean(file_path, on_page_progress=on_page_progress)
